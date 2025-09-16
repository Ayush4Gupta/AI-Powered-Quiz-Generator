# document_extractors.py
"""
Document text extraction utilities for multiple file formats.
Follows the same pattern as the existing PDF extractor in pdf.py
"""
from pathlib import Path
from typing import List, Dict, Any
import os

class DocumentExtractionError(Exception):
    """Raised when document extraction fails"""
    pass

def extract_text_from_docx(docx_path: Path) -> str:
    """Extract plain text from DOCX files using python-docx"""
    try:
        from docx import Document
        doc = Document(str(docx_path))
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        return "\n".join(text_parts)
    except ImportError:
        raise DocumentExtractionError("python-docx not installed. Install with: pip install python-docx")
    except Exception as exc:
        raise DocumentExtractionError(f"Could not parse DOCX: {exc}") from exc

def extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract plain text from PPTX files using python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = shape.text.strip()
                    if shape_text:
                        slide_texts.append(shape_text)
            
            if slide_texts:
                text_parts.append(f"=== Slide {slide_num} ===")
                text_parts.extend(slide_texts)
                text_parts.append("")  # Empty line between slides
        
        return "\n".join(text_parts)
    except ImportError:
        raise DocumentExtractionError("python-pptx not installed. Install with: pip install python-pptx")
    except Exception as exc:
        raise DocumentExtractionError(f"Could not parse PPTX: {exc}") from exc

def extract_text_from_txt(txt_path: Path) -> str:
    """Extract plain text from TXT files"""
    try:
        with open(txt_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encodings
        try:
            with open(txt_path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception as exc:
            raise DocumentExtractionError(f"Could not read TXT file with UTF-8 or Latin-1 encoding: {exc}") from exc
    except Exception as exc:
        raise DocumentExtractionError(f"Could not parse TXT: {exc}") from exc

def extract_text_from_pdf(pdf_path: Path) -> str:
    """Extract plain text from PDF files using existing PDF utility"""
    try:
        # Use the existing PDF extractor
        from app.utils.pdf import extract_text
        return extract_text(pdf_path)
    except Exception as exc:
        raise DocumentExtractionError(f"Could not parse PDF: {exc}") from exc

# Supported file types and their extractors
SUPPORTED_EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".txt": extract_text_from_txt,
}

def get_supported_extensions() -> List[str]:
    """Get list of supported file extensions"""
    return list(SUPPORTED_EXTRACTORS.keys())

def is_supported_file(file_path: Path) -> bool:
    """Check if file extension is supported"""
    return file_path.suffix.lower() in SUPPORTED_EXTRACTORS

def extract_text_from_file(file_path: Path) -> str:
    """Extract text from any supported file type"""
    file_path = Path(file_path)
    
    if not file_path.exists():
        raise DocumentExtractionError(f"File not found: {file_path}")
    
    if not file_path.is_file():
        raise DocumentExtractionError(f"Path is not a file: {file_path}")
    
    extension = file_path.suffix.lower()
    
    if extension not in SUPPORTED_EXTRACTORS:
        raise DocumentExtractionError(
            f"Unsupported file type: {extension}. "
            f"Supported types: {', '.join(get_supported_extensions())}"
        )
    
    extractor = SUPPORTED_EXTRACTORS[extension]
    return extractor(file_path)
